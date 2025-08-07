# mcp_server_tools.py
from agents.mcp import MCPToolServer, tool
import os, fitz
from pptx import Presentation

class FileOpsMCPServer(MCPToolServer):
    @tool
    def list_files(self, source="local"):
        if source == "local":
            files = [f for f in os.listdir(os.getenv("LOCAL_FILE_DIR"))
                     if f.lower().endswith(('.pdf', '.ppt', '.pptx'))]
            return files
        # Placeholder: Add Azure listing

    @tool
    def extract_pdf(self, filename: str):
        path = os.path.join(os.getenv("LOCAL_FILE_DIR"), filename)
        doc = fitz.open(path)
        return "\n".join([page.get_text() for page in doc])

    @tool
    def extract_ppt(self, filename: str):
        path = os.path.join(os.getenv("LOCAL_FILE_DIR"), filename)
        prs = Presentation(path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

if __name__ == "__main__":
    FileOpsMCPServer().serve()
