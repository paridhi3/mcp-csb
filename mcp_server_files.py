"""
MCP Tool Server for file listing and text extraction from PDFs and PPTs.
This server exposes tools to list files and extract text content.
"""

import os
import fitz  # PyMuPDF
from pptx import Presentation
from dotenv import load_dotenv
from agents.mcp import MCPToolServer, tool

load_dotenv()

LOCAL_DIR = os.getenv("LOCAL_FILE_DIR", "./case_studies")


class FileOpsMCPServer(MCPToolServer):
    @tool
    def list_files(self, source: str = "local"):
        """
        Lists available files from local directory.
        Supports PDFs, PPT, PPTX files.
        """
        if source == "local":
            files = [f for f in os.listdir(LOCAL_DIR)
                     if f.lower().endswith((".pdf", ".ppt", ".pptx"))]
            return files
        else:
            # Placeholder for Azure blob or other sources
            return []

    @tool
    def extract_pdf(self, filename: str):
        """
        Extract text content from a PDF file.
        """
        file_path = os.path.join(LOCAL_DIR, filename)
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    @tool
    def extract_ppt(self, filename: str):
        """
        Extract text content from a PPT or PPTX file.
        """
        file_path = os.path.join(LOCAL_DIR, filename)
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
        return "\n".join(text_runs)


if __name__ == "__main__":
    # Start MCP server for file ops
    FileOpsMCPServer().serve()
